# -*- coding: utf-8 -*-

import os
import logging
import base64
import asyncio
import httpx
from telegram import Update, InlineKeyboardMarkup, InlineKeyboardButton
from telegram.ext import Application, CommandHandler, MessageHandler, filters, ContextTypes, CallbackQueryHandler, PreCheckoutQueryHandler
from dotenv import load_dotenv
from telegram.error import RetryAfter, NetworkError, BadRequest
from PIL import Image
import io
import aiohttp
from collections import defaultdict
import json
import tempfile
import pathlib

# ВАЖНО: Для работы с презентациями необходимо установить библиотеку python-pptx:
# pip install python-pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt # Добавлено для работы с размерами
from pptx.enum.dml import MSO_THEME_COLOR # Добавлено для работы с цветами темы
from pptx.enum.shapes import MSO_SHAPE # Добавлено для работы с фигурами
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE # Добавлено для настройки текста

# --- Настройка и конфигурация ---

# Загружаем переменные из файла secrets.env
load_dotenv('secrets.env')

# Настройка логирования
# Уровень логирования можно изменить здесь:
# logging.INFO (стандартный), logging.DEBUG (очень подробный)
logging.basicConfig(
    format="%(asctime)s - %(name)s - %(levelname)s - %(message)s", level=logging.INFO
)
logging.getLogger("aiohttp.internal_client").setLevel(logging.WARNING)
LOGGER = logging.getLogger(__name__)

# Получение токена бота и ключей API
BOT_TOKEN = os.getenv("BOT_TOKEN") # Используем переменную окружения
GEMINI_API_KEYS = os.getenv("GEMINI_API_KEYS", "").split(',')

# Проверка, что API-ключи существуют
if not all(GEMINI_API_KEYS):
    LOGGER.error("Ошибка: API-ключи не найдены в файле secrets.env. Пожалуйста, убедитесь, что они там есть.")
    exit(1)

# Переменная для отслеживания текущего индекса ключа
key_index = 0

# Максимальное количество попыток для запроса к API
MAX_RETRIES = 5
RETRY_DELAY = 1

# Максимальное количество сообщений в истории диалога для каждого пользователя
MAX_HISTORY_MESSAGES = 10

# Словари для хранения данных по пользователям
user_history = {}
# Единый словарь для всех настроек пользователя, включая кредиты.
# Это гарантирует, что все функции будут обращаться к одному и тому же состоянию.
user_settings = defaultdict(lambda: {"response_format": "html", "html_credits": 0})
media_groups = {} # Словарь для временного хранения медиагрупп
# ВАЖНО: Переключатель для тестового режима. Установи в False для реальных платежей.
IS_TEST_MODE = True

# CSS-стили для тетрадей
NOTEBOOK_STYLES = """
<style>
    body {
        font-family: 'Times New Roman', Times, serif;
        margin: 0;
        padding: 0;
    }

    .math-background {
        max-width: 35ch;
        height: 100%;
        background-color: #ffffff;
        background-image:
            linear-gradient(to right, #add8e6 0.5px, transparent 0.5px),
            linear-gradient(to bottom, #add8e6 0.5px, transparent 0.5px);
        background-size: 5mm 5mm;
        padding: 10mm;
        box-sizing: border-box;
    }

    .russian-background {
        max-width: 35ch;
        height: 100%;
        background-color: #ffffff;
        background-image: repeating-linear-gradient(to bottom, #add8e6 0, #add8e6 1px, transparent 1px, transparent 10mm);
        background-size: 100% 10mm;
        padding: 10mm;
        box-sizing: border-box;
    }

    .default-background {
        background-color: #f0f0f0;
        padding: 20px;
        font-family: Arial, sans-serif;
    }

    h1, h2, h3, h4, h5, h6, p {
        margin: 0;
        line-height: 1.5;
    }

    .diagram {
        border: 2px solid #3498db;
        border-radius: 8px;
        padding: 10px;
        margin: 10px 0;
    }
</style>
"""

# Единый промпт для всех запросов
DEVELOPER_PROMPT = f"""
Ты — продвинутая нейросеть, способная отвечать на вопросы, анализировать изображения и генерировать HTML-код. Твоя основная задача — быть полезным и точным ассистентом.
Твои ответы должны быть лаконичными и по существу, отвечай на русском языке. Будь лаконичен в своих ответах. Избегай чрезмерно длинного текста.
**ВАЖНО:** Ты должен генерировать **ТОЛЬКО** HTML-код. Весь твой ответ должен начинаться с `<!DOCTYPE html>` и заканчиваться `</html>`. Не используй никаких других символов или текста вне HTML-тегов.

**Инструкции:**
- **Ограничение по размеру:** В твоём ответе ширина контента не должна превышать 35 "клеток" (символов) в ширину. Для этого используй CSS-стили, чтобы ограничить ширину основного контейнера (например, 'max-width: 35ch;'). Вертикальная высота не ограничена.
- **Оформление:** Для диаграмм и чертежей используй красивые, яркие цвета и линии. Можешь выделять важный текст разными цветами.
- **Подробность:** Предоставляй оптимальный по объёму, полный и понятный ответ. Пиши всё, что необходимо для решения, без "ленивых" сокращений.
- **Диаграммы:** Если задание требует диаграммы или чертежа, создай его в HTML. Диаграмма должна быть красивой, понятной и полностью соответствовать заданию, а не быть просто примером.
- **Использование фона:** Используй фон тетрадного листа (.math-background или .russian-background) только для школьных заданий. Если задание не является школьным, используй простой фон (.default-background).
- Делай дизайн красивым, современным и адаптивным. Убедись, что макет хорошо выглядит и работает как на мобильных устройствах в портретной и альбомной ориентации, так и на компьютерах.
- Помни, что ты создаешь готовую домашнюю работу, которую мог бы написать ученик.
- **Для фото:** Если тебе присылают фотографию, **всегда** предполагай, что это задание, которое нужно **решить**, а не просто описать. Используй подпись к фото для уточнения контекста, но основной задачей всегда является решение.
- **Для файлов:** Если тебе присылают файлы (HTML, Python, PDF, TXT и т.д.), анализируй их содержимое и давай осмысленный ответ, основанный на этом содержимом. Если есть подпись, используй ее как дополнительный контекст.

Пример для школьного задания:
<!DOCTYPE html>
<html lang="ru">
<head>
    <title>Решение</title>
    {NOTEBOOK_STYLES}
</head>
<body>
    <div class="math-background">
        <h1>Задание по алгебре</h1>
        <p><strong>Решение:</strong></p>
        <p>2x + 5 = 11</p>
        <p>2x = 6</p>
        <p>x = 3</p>
    </div>
</body>
</html>

Пример для обычного запроса:
<!DOCTYPE html>
<html lang="ru">
<head>
    <title>Ответ</title>
    {NOTEBOOK_STYLES}
</head>
<body>
    <div class="default-background">
        <h1>Информация о Python</h1>
        <p>Python — это высокоуровневый интерпретируемый язык программирования...</p>
    </div>
</body>
</html>
"""

# Промпт для режима презентации
PRESENTATION_PROMPT = """
You are an advanced AI designed to create structured JSON output for generating PowerPoint presentations.
Your task is to analyze the user's request and generate a JSON object that will be used to create the slides.

**IMPORTANT:** You must generate **ONLY** JSON. No other characters or text.

**Instructions:**
- The JSON should be an array of objects.
- Each object in the array represents a single slide.
- Each object must have two keys:
    - `"title"`: A string for the slide's title.
    - `"points"`: An array of strings with the main points or bullet points for that slide.

**Example of the expected JSON structure:**
[
  {
    "title": "Slide 1 Title",
    "points": [
      "Point 1",
      "Point 2",
      "Point 3"
    ]
  },
  {
    "title": "Slide 2 Title",
    "points": [
      "Another point",
      "Yet another point"
    ]
  }
]
"""
# --- Функции API и вспомогательные функции ---

async def get_next_api_key():
    """Возвращает следующий API-ключ из списка по кругу."""
    global key_index
    api_key = GEMINI_API_KEYS[key_index]
    key_index = (key_index + 1) % len(GEMINI_API_KEYS)
    LOGGER.info(f"Переключение на API ключ с индексом {key_index - 1}.")
    return api_key

async def call_gemini_api(payload: dict) -> str:
    """Отправляет запрос к Gemini API и возвращает ответ."""
    async with aiohttp.ClientSession() as session:
        retries = 0
        while retries < MAX_RETRIES:
            api_key = await get_next_api_key()
            api_url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-2.5-flash-preview-05-20:generateContent?key={api_key}"
            LOGGER.info(f"Попытка {retries + 1}/{MAX_RETRIES} с API ключом {key_index}.")
            LOGGER.debug(f"Payload для Gemini API: {json.dumps(payload, indent=2)}")
            try:
                async with session.post(api_url, json=payload, timeout=300.0) as response:
                    LOGGER.info(f"Ответ от Gemini API: HTTP {response.status}")
                    if response.status == 400:
                        error_text = await response.text()
                        if "API key not valid" in error_text:
                            LOGGER.error(f"Неверный API ключ: {api_key}. Переключаюсь на следующий.")
                            retries += 1
                            continue
                        else:
                            LOGGER.error(f"HTTP 400 Bad Request. Details: {error_text}")
                            return "Извините, этот тип файла не поддерживается или запрос неверно сформирован."
                    
                    response.raise_for_status()
                    result = await response.json()
                    text_content = result.get('candidates', [{}])[0].get('content', {}).get('parts', [{}])[0].get('text', 'Не удалось получить ответ.')
                    LOGGER.info("Успешный ответ от Gemini API.")
                    return text_content
            except aiohttp.ClientResponseError as e:
                LOGGER.error(f"HTTP error during Gemini API request: {e.status} - {e.message}")
                if e.status == 429:
                    LOGGER.warning("Rate limit exceeded for Gemini API. Switching to next key...")
                    retries += 1
                else:
                    retries += 1
                    await asyncio.sleep(RETRY_DELAY * (2 ** retries))
            except aiohttp.ClientError as e:
                LOGGER.error(f"Network error during Gemini API request: {e}")
                retries += 1
                await asyncio.sleep(RETRY_DELAY * (2 ** retries))
            except Exception as e:
                LOGGER.error(f"Unknown error: {e}")
                return "Извините, произошла ошибка."
        LOGGER.error("Не удалось получить ответ от нейросети после нескольких попыток.")
        return "Извините, не удалось получить ответ от нейросети после нескольких попыток."

async def send_html_file(update: Update, html_code: str):
    """Creates and sends an HTML file from the generated code."""
    LOGGER.info("Начало отправки HTML-файла.")
    try:
        file_size_bytes = len(html_code.encode('utf-8'))
        LOGGER.info(f"Размер генерируемого HTML-файла: {file_size_bytes} байт.")
        if file_size_bytes > 50 * 1024 * 1024:
            await update.message.reply_text("Извините, сгенерированный файл слишком большой для отправки в Telegram.")
            LOGGER.warning(f"Файл слишком большой для отправки: {file_size_bytes} байт.")
            return

        await update.message.reply_document(
            document=html_code.encode('utf-8'),
            filename="solution.html"
        )
        LOGGER.info(f"HTML-файл успешно отправлен пользователю {update.effective_user.id}")
    except BadRequest as e:
        if "too long" in str(e).lower():
            await update.message.reply_text("Извините, сгенерированный файл слишком большой для отправки в Telegram. Пожалуйста, попробуйте сформулировать запрос более кратко.")
        else:
            LOGGER.error(f"Ошибка при отправке HTML-файла: {e}")
            await update.message.reply_text("Извините, произошла ошибка при отправке файла.")
    except Exception as e:
        LOGGER.error(f"Ошибка при отправке HTML-файла: {e}")
        await update.message.reply_text("Извините, произошла ошибка при отправке файла.")

async def create_and_send_pptx_file(update: Update, slides_data: list):
    """
    Создает PowerPoint-презентацию из JSON-данных с улучшенным дизайном.
    Использует новые настройки для цветов, шрифтов и макета.
    """
    LOGGER.info("Начало создания и отправки PPTX-файла.")
    prs = Presentation()

    # Настройка цветов и стилей
    BACKGROUND_COLOR = RGBColor(245, 245, 245)
    PRIMARY_COLOR = RGBColor(41, 128, 185) # Синий
    SECONDARY_COLOR = RGBColor(52, 73, 94) # Темно-серый
    ACCENT_COLOR = RGBColor(52, 152, 219) # Голубой

    for slide_info in slides_data:
        LOGGER.debug(f"Обработка слайда: {slide_info.get('title', 'Без заголовка')}")
        # Выбор макета для слайда с заголовком и списком
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # Настройка фона слайда
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = BACKGROUND_COLOR

        # Добавление фигуры-акцента вверху слайда
        left = top = Inches(0)
        width = prs.slide_width
        height = Inches(0.2)
        accent_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        accent_shape.fill.solid()
        accent_shape.fill.fore_color.rgb = ACCENT_COLOR
        accent_shape.line.fill.background()
        
        # Добавление тени для акцентной фигуры
        accent_shape.shadow.visible = True
        accent_shape.shadow.offset_x = Inches(0)
        accent_shape.shadow.offset_y = Inches(0.1)

        # Работа с заголовком
        title_shape = slide.shapes.title
        title_shape.text = slide_info.get("title", "Без заголовка")
        title_tf = title_shape.text_frame
        
        # Настройка шрифта и выравнивания заголовка
        title_para = title_tf.paragraphs[0]
        title_para.font.name = 'Arial'
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = PRIMARY_COLOR
        title_para.alignment = PP_ALIGN.CENTER
        
        # Улучшенное позиционирование и размер заголовка
        title_shape.top = Inches(0.5)
        title_shape.height = Inches(1.5)
        title_shape.width = prs.slide_width - Inches(2)
        title_shape.left = Inches(1)

        # Работа с основным текстом
        if "points" in slide_info and isinstance(slide_info["points"], list):
            content_shape = slide.placeholders[1]
            content_tf = content_shape.text_frame
            content_tf.clear()
            content_tf.word_wrap = True
            
            for point in slide_info["points"]:
                p = content_tf.add_paragraph()
                p.text = point
                p.level = 0
                p.font.name = 'Arial'
                p.font.size = Pt(20)
                p.font.color.rgb = SECONDARY_COLOR
                p.alignment = PP_ALIGN.LEFT
                p.space_after = Pt(10)
        
    # Сохранение во временный файл
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as temp_file:
        prs.save(temp_file.name)
        filepath = temp_file.name
    
    # Отправка файла
    try:
        await update.message.reply_document(
            document=open(filepath, 'rb'),
            filename="presentation.pptx"
        )
        LOGGER.info(f"Presentation file sent successfully to {update.effective_user.id}")
    except Exception as e:
        LOGGER.error(f"Error sending PowerPoint file: {e}")
        await update.message.reply_text("Извините, произошла ошибка при отправке файла презентации.")
    finally:
        # Удаление временного файла
        os.remove(filepath)
        LOGGER.info(f"Временный файл {filepath} удален.")

# --- Обработчики команд и сообщений ---

# Обработчик команды /start
async def start_command_handler(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Отправляет приветственное сообщение с меню в виде кнопок."""
    user_id = update.effective_user.id
    LOGGER.info(f"Пользователь {user_id} отправил команду /start.")
    if user_id in user_history:
        del user_history[user_id]
    
    welcome_message = """
Привет! Я - твой личный помощник для учебы.
Я могу решить твои домашние задания по разным предметам.

⭐ **Как пользоваться:**
Просто отправь мне фотографию или файл (PDF, TXT, HTML, PY и т.д.) с заданием. Ты можешь добавить подпись к фото или файлу, чтобы дать мне подсказку.

⚠️ **Важно для геометрии!**
Для лучшего результата отправляй геометрические задачи **по одной**, а не несколько на одном листе. Это поможет мне избежать путаницы и дать более точное решение.

Я всегда буду отвечать тебе красивым HTML-файлом, оформленным как тетрадный лист.
"""
    # Создаем кнопки
    keyboard = [
        [InlineKeyboardButton("Начать чат", callback_data='start_chat')],
        [InlineKeyboardButton("Настройки", callback_data='settings')],
        [InlineKeyboardButton("Донат", callback_data='donate')] # Изменено на новый обработчик
    ]
    reply_markup = InlineKeyboardMarkup(keyboard)

    await update.message.reply_text(welcome_message, reply_markup=reply_markup)
    LOGGER.info(f"Пользователь {user_id} начал чат.")

# НОВОЕ: Обработчик команды /donate
async def donate_command_handler(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Отправляет меню для доната."""
    user_id = update.effective_user.id
    LOGGER.info(f"Пользователь {user_id} отправил команду /donate.")
    
    message = "Спасибо за вашу поддержку! Выберите, сколько звёзд вы хотите купить. За каждую звезду вы получите 10 ответов."
    keyboard = [
        [InlineKeyboardButton("1 ⭐️", callback_data='buy_stars_1')],
        [InlineKeyboardButton("5 ⭐️", callback_data='buy_stars_5')],
        [InlineKeyboardButton("10 ⭐️", callback_data='buy_stars_10')],
        [InlineKeyboardButton("Назад", callback_data='settings')]
    ]
    reply_markup = InlineKeyboardMarkup(keyboard)
    
    if update.callback_query:
        await update.callback_query.edit_message_text(message, reply_markup=reply_markup)
    else:
        await update.message.reply_text(message, reply_markup=reply_markup)
    LOGGER.info(f"Меню доната отправлено пользователю {user_id}.")


# Обработчик команды /reset
async def reset_command_handler(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Очищает историю диалога для текущего пользователя."""
    user_id = update.effective_user.id
    LOGGER.info(f"Пользователь {user_id} отправил команду /reset.")
    if user_id in user_history:
        del user_history[user_id]
        LOGGER.info(f"История диалога для пользователя {user_id} очищена.")
    await update.message.reply_text("Диалог сброшен. Можете начинать новую беседу.")

# НОВАЯ функция для тестирования: выдает кредиты
async def get_stars_handler(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Добавляет 5 кредитов для тестирования."""
    user_id = update.effective_user.id
    LOGGER.info(f"Пользователь {user_id} отправил команду /get_stars.")
    if IS_TEST_MODE:
        user_settings[user_id]["html_credits"] += 5
        LOGGER.info(f"Добавлено 5 кредитов пользователю {user_id}. Текущий баланс: {user_settings[user_id]['html_credits']}")
        await update.message.reply_text(
            f"✅ **Режим тестирования:** Вам добавлено 5 кредитов. Теперь вы можете генерировать HTML-ответы."
            f"\n\n**Текущий баланс:** {user_settings[user_id]['html_credits']} кредитов."
        )
    else:
        LOGGER.info("Команда /get_stars была вызвана в не-тестовом режиме.")
        await update.message.reply_text("Эта команда доступна только в режиме тестирования.")

# Обработчик кнопок
async def button_handler(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Обрабатывает нажатия на кнопки в меню."""
    query = update.callback_query
    await query.answer()
    user_id = query.from_user.id
    LOGGER.info(f"Пользователь {user_id} нажал кнопку: {query.data}")

    if query.data == 'start_chat':
        await query.edit_message_text("Отлично, можете присылать ваши задания. Чтобы начать заново, используйте команду /start.")
    elif query.data == 'settings':
        keyboard = [
            [InlineKeyboardButton("Способ отправки", callback_data='settings_send_method')],
            [InlineKeyboardButton("Купить ответы (10 за 1⭐️)", callback_data='donate')], # Изменено, ведёт в новое меню доната
            [InlineKeyboardButton("Назад", callback_data='start_chat')]
        ]
        reply_markup = InlineKeyboardMarkup(keyboard)
        await query.edit_message_text("Выберите, как мне отвечать:", reply_markup=reply_markup)
    elif query.data == 'donate':
        await donate_command_handler(update, context) # Новый вызов
    elif query.data.startswith('buy_stars_'):
        stars_cost = int(query.data.split('_')[2])
        LOGGER.info(f"Пользователь {user_id} инициировал покупку {stars_cost} звёзд.")
        await send_invoice(update, context, stars_cost)
    elif query.data == 'settings_send_method':
        keyboard = [
            [InlineKeyboardButton("HTML (файл)", callback_data='format_html')],
            [InlineKeyboardButton("Презентация (файл)", callback_data='format_presentation')],
            [InlineKeyboardButton("Текст (сообщение)", callback_data='format_text')],
            [InlineKeyboardButton("Назад", callback_data='settings')]
        ]
        reply_markup = InlineKeyboardMarkup(keyboard)
        await query.edit_message_text("Выберите формат ответа:", reply_markup=reply_markup)
    elif query.data.startswith('format_'):
        response_format = query.data.split('_')[1]
        user_settings[user_id]["response_format"] = response_format
        LOGGER.info(f"Формат ответа для пользователя {user_id} изменен на: {response_format}")
        await query.edit_message_text(f"Отлично! Теперь я буду отвечать в формате **{response_format.upper()}**. Чтобы изменить, зайдите в /settings.")

async def process_media_group(media_group_id, user_id, context):
    """
    Ждет небольшое время, чтобы убедиться, что все сообщения из медиагруппы получены,
    затем обрабатывает их.
    """
    LOGGER.info(f"Получена медиагруппа с ID: {media_group_id} от пользователя {user_id}. Ожидание...")
    await asyncio.sleep(2) # Ждем 2 секунды, чтобы собрать все фото
    
    if media_group_id not in media_groups:
        LOGGER.warning(f"Медиагруппа {media_group_id} уже была обработана или не найдена.")
        return # Если медиагруппа уже обработана, выходим

    messages = media_groups.pop(media_group_id)["messages"]
    LOGGER.info(f"Собрано {len(messages)} сообщений из медиагруппы {media_group_id}. Начало обработки.")
    
    await context.bot.send_message(user_id, "⌛ Обрабатываю ваш альбом...")
    
    content_parts = []
    caption = ""
    
    for msg in messages:
        if msg.photo:
            photo_file = await msg.photo[-1].get_file()
            photo_content = io.BytesIO(await photo_file.download_as_bytearray())
            try:
                img = Image.open(photo_content)
                img_rgb = img.convert("RGB")
                buffer = io.BytesIO()
                img_rgb.save(buffer, format="JPEG")
                base64_encoded_image = base64.b64encode(buffer.getvalue()).decode('utf-8')
                
                content_parts.append({
                    "inlineData": {
                        "mimeType": "image/jpeg",
                        "data": base64_encoded_image
                    }
                })
            except Exception as e:
                LOGGER.error(f"Ошибка при обработке изображения из медиагруппы: {e}")
                continue
            
            if msg.caption:
                caption = msg.caption
    
    text_prompt = caption or "Реши эти задания."
    
    payload = {
        "contents": [
            {
                "parts": [
                    {"text": DEVELOPER_PROMPT + "\n\n" + text_prompt},
                ] + content_parts
            }
        ],
        "generationConfig": {"temperature": 0.4}
    }
    
    try:
        html_response = await call_gemini_api(payload)
        await send_html_file(messages[0], html_response)
    except Exception as e:
        LOGGER.error(f"Ошибка при обработке медиагруппы: {e}")
        await context.bot.send_message(user_id, "Извините, произошла ошибка при обработке альбома.")

# Изменено: теперь принимает stars_cost
async def send_invoice(update: Update, context: ContextTypes.DEFAULT_TYPE, stars_cost: int) -> None:
    """Создаёт и отправляет счёт для оплаты Telegram Stars."""
    user_id = update.effective_user.id
    
    invoice_payload = f"html_purchase_{stars_cost}" # Уникальный идентификатор платежа с количеством звёзд
    title = f"Покупка {stars_cost} звёзд"
    description = f"Получите {stars_cost * 10} ответов в формате HTML."
    
    try:
        await context.bot.send_invoice(
            chat_id=user_id,
            title=title,
            description=description,
            payload=invoice_payload,
            provider_token="", # Оставляем пустым, так как это Telegram Stars
            currency="XTR", # "XTR" - код валюты для Telegram Stars
            prices=[{"label": f"{stars_cost} звезда" if stars_cost == 1 else f"{stars_cost} звёзд", "amount": stars_cost}],
            start_parameter="html_purchase", # Уникальный параметр
            need_name=False,
            need_shipping_address=False
        )
        LOGGER.info(f"Счёт для оплаты {stars_cost} звёзд отправлен пользователю {user_id}")
    except Exception as e:
        LOGGER.error(f"Ошибка при отправке счёта: {e}")
        await context.bot.send_message(user_id, "Извините, не удалось создать счёт для оплаты.")

# НОВЫЙ обработчик для pre_checkout_query
async def pre_checkout_query_handler(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Обрабатывает pre-checkout запросы от Telegram."""
    query = update.pre_checkout_query
    LOGGER.info(f"Получен pre_checkout_query от {query.from_user.id}. Payload: {query.invoice_payload}")

    # Проверяем, что payload начинается с нужного префикса
    if query.invoice_payload.startswith("html_purchase_"):
        await query.answer(ok=True)
        LOGGER.info(f"Pre_checkout_query успешно подтвержден для пользователя {query.from_user.id}")
    else:
        await query.answer(ok=False, error_message="Извините, мы не можем обработать этот платёж.")
        LOGGER.warning(f"Неизвестный payload в pre_checkout_query от {query.from_user.id}")

async def successful_payment_callback(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Обрабатывает успешную оплату и выдаёт "кредиты"."""
    user_id = update.effective_user.id
    payload = update.effective_message.successful_payment.invoice_payload
    
    LOGGER.info(f"Получен успешный платёж от пользователя {user_id}. Payload: {payload}")
    
    if payload.startswith("html_purchase_"):
        try:
            stars_bought = int(payload.split('_')[2])
            credits_to_add = stars_bought * 10 # 10 ответов за каждую звезду
            
            # Увеличиваем "баланс" пользователя
            user_settings[user_id]["html_credits"] += credits_to_add
            LOGGER.info(f"Баланс пользователя {user_id} увеличен на {credits_to_add}. Текущий баланс: {user_settings[user_id]['html_credits']}")
            
            await update.effective_message.reply_text(f"✅ Оплата прошла успешно! Вам зачислено {credits_to_add} кредитов. Теперь вы можете получить ответы в формате HTML. Пожалуйста, отправьте мне ваше задание.")
        except (IndexError, ValueError) as e:
            LOGGER.error(f"Не удалось распарсить payload {payload}: {e}")
            await update.effective_message.reply_text("Извините, произошла ошибка при зачислении кредитов. Пожалуйста, свяжитесь с администратором.")

async def handle_message(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Универсальный обработчик для текста, фото и файлов."""
    
    user_id = update.effective_user.id
    LOGGER.info(f"Получено сообщение от пользователя {user_id}.")
    
    response_format = user_settings[user_id].get("response_format", "html")
    
    if update.message.media_group_id:
        media_group_id = update.message.media_group_id
        
        if media_group_id not in media_groups:
            media_groups[media_group_id] = {"messages": [], "task": None}
        
        media_groups[media_group_id]["messages"].append(update.message)
        
        if media_groups[media_group_id]["task"]:
            media_groups[media_group_id]["task"].cancel()
            
        media_groups[media_group_id]["task"] = asyncio.create_task(
            process_media_group(media_group_id, user_id, context)
        )
        return

    # Проверяем, нужно ли обрабатывать как HTML-файл и есть ли "кредиты"
    if response_format == "html":
        credits = user_settings[user_id].get("html_credits", 0)
        LOGGER.info(f"Пользователь {user_id} запросил HTML-формат. Текущий баланс: {credits}")
        
        # Этот блок кода проверяет, достаточно ли кредитов у пользователя для получения ответа.
        # Если у пользователя 1 или более кредит, условие "credits <= 0" будет ложным,
        # и бот продолжит работу. Если кредитов 0, бот попросит оплату.
        if credits <= 0:
            LOGGER.info(f"У пользователя {user_id} недостаточно кредитов. Отправка предложения о покупке.")
            await update.message.reply_text(
                "Чтобы получить ответ в формате HTML, у вас должен быть как минимум 1 кредит. Вы можете купить их в меню: /donate.",
                reply_markup=InlineKeyboardMarkup([
                    [InlineKeyboardButton("Перейти в меню доната", callback_data='donate')]
                ])
            )
            return
        
        # Если "кредит" есть, уменьшаем его на 1 перед началом генерации
        user_settings[user_id]["html_credits"] -= 1
        LOGGER.info(f"1 кредит использован. Новый баланс для {user_id}: {user_settings[user_id]['html_credits']}")
        await update.message.reply_text(f"⏳ Использую 1 «кредит». Осталось: {user_settings[user_id]['html_credits']}. Обрабатываю ваш запрос...")
    else:
        LOGGER.info(f"Пользователь {user_id} запросил формат: {response_format}. Обрабатываю запрос.")
        await update.message.reply_text("⏳ Обрабатываю ваш запрос...")
        
    if user_id not in user_history:
        user_history[user_id] = []
        LOGGER.debug(f"Инициализирована новая история для пользователя {user_id}.")

    while len(user_history[user_id]) >= MAX_HISTORY_MESSAGES:
        user_history[user_id].pop(0)
        LOGGER.debug("Удалено самое старое сообщение из истории диалога.")

    contents = [{ "role": "user", "parts": [{ "text": DEVELOPER_PROMPT }]}]
    for item in user_history[user_id]:
        contents.append(item)

    if update.message.document:
        document = update.message.document
        LOGGER.info(f"Получен документ от {user_id}: {document.file_name}, MIME-тип: {document.mime_type}")
        file_info = {
            'file_id': document.file_id,
            'file_name': document.file_name,
            'mime_type': document.mime_type
        }
        file = await context.bot.get_file(file_info['file_id'])
        file_content = await file.download_as_bytearray()
        
        if file_info['mime_type'].startswith('image/'):
            LOGGER.info(f"Документ идентифицирован как изображение. Размер: {len(file_content)} байт.")
            try:
                img = Image.open(io.BytesIO(file_content))
                img_rgb = img.convert("RGB")
                buffer = io.BytesIO()
                img_rgb.save(buffer, format="JPEG")
                base64_encoded_image = base64.b64encode(buffer.getvalue()).decode('utf-8')
                
                text_prompt = update.message.caption if update.message.caption else "Проанализируй это изображение."
                contents.append({
                    "role": "user",
                    "parts": [
                        {"text": text_prompt},
                        {
                            "inlineData": {
                                "mimeType": "image/jpeg",
                                "data": base64_encoded_image
                            }
                        }
                    ]
                })
            except Exception as e:
                LOGGER.error(f"Ошибка при обработке изображения из документа: {e}")
                await update.message.reply_text("Извините, произошла ошибка при обработке изображения.")
                return
        
        elif file_info['mime_type'].startswith('text/') or file_info['file_name'].lower().endswith(('.py', '.txt', '.html', '.md')):
            LOGGER.info(f"Документ идентифицирован как текстовый файл.")
            try:
                decoded_content = file_content.decode('utf-8')
                text_prompt = update.message.caption if update.message.caption else ""
                contents.append({
                    "role": "user",
                    "parts": [{ "text": f"{text_prompt}\n\nСодержимое файла:\n\n{decoded_content}"}]
                })
                LOGGER.debug("Содержимое текстового файла успешно декодировано.")
            except UnicodeDecodeError:
                LOGGER.error("Не удалось декодировать файл. Возможно, это бинарный файл.")
                await update.message.reply_text("Извините, не удалось прочитать этот файл как текст.")
                return
        else:
            LOGGER.warning(f"Получен неподдерживаемый тип файла: {file_info['mime_type']}.")
            await update.message.reply_text("Извините, этот тип файла не поддерживается.")
            return
            
    elif update.message.photo:
        photo = update.message.photo[-1]
        LOGGER.info(f"Получена фотография от {user_id}. File ID: {photo.file_id}")
        file = await context.bot.get_file(photo.file_id)
        file_content = await file.download_as_bytearray()
        
        try:
            img = Image.open(io.BytesIO(file_content))
            img_rgb = img.convert("RGB")
            buffer = io.BytesIO()
            img_rgb.save(buffer, format="JPEG")
            base64_encoded_image = base64.b64encode(buffer.getvalue()).decode('utf-8')
            
            text_prompt = update.message.caption if update.message.caption else "Проанализируй это изображение."
            contents.append({
                "role": "user",
                "parts": [
                    {"text": text_prompt},
                    {
                        "inlineData": {
                            "mimeType": "image/jpeg",
                            "data": base64_encoded_image
                        }
                    }
                ]
            })
            LOGGER.info("Фотография успешно обработана и добавлена в запрос.")
        except Exception as e:
            LOGGER.error(f"Ошибка при обработке фотографии: {e}")
            await update.message.reply_text("Извините, произошла ошибка при обработке фотографии.")
            return

    elif update.message.text:
        text_prompt = update.message.text
        LOGGER.info(f"Получен текстовый запрос от {user_id}: '{text_prompt}'")
        contents.append({
            "role": "user",
            "parts": [{"text": text_prompt}]
        })
    else:
        await update.message.reply_text("Пожалуйста, предоставьте текст, фотографию или файл, чтобы я мог помочь.")
        return

    try:
        if response_format == "html":
            payload = {
                "contents": contents,
                "generationConfig": {"temperature": 0.4}
            }
            LOGGER.info("Отправка запроса в Gemini API для генерации HTML.")
            gemini_response = await call_gemini_api(payload)
            
            user_history[user_id].append({
                "role": "model",
                "parts": [{"text": gemini_response}]
            })
            
            await send_html_file(update, gemini_response)

        elif response_format == "presentation":
            LOGGER.info("Отправка запроса в Gemini API для генерации презентации (JSON).")
            # Change the prompt and force JSON output for presentation mode
            contents[0]["parts"][0]["text"] = PRESENTATION_PROMPT
            
            payload = {
                "contents": contents,
                "generationConfig": {
                    "temperature": 0.4,
                    "responseMimeType": "application/json",
                    "responseSchema": {
                        "type": "ARRAY",
                        "items": {
                            "type": "OBJECT",
                            "properties": {
                                "title": {"type": "STRING"},
                                "points": {
                                    "type": "ARRAY",
                                    "items": {"type": "STRING"}
                                }
                            },
                            "propertyOrdering": ["title", "points"]
                        }
                    }
                }
            }
            
            gemini_json_response = await call_gemini_api(payload)

            try:
                slides_data = json.loads(gemini_json_response)
                LOGGER.info("JSON для презентации успешно разобран.")
                await create_and_send_pptx_file(update, slides_data)
            except json.JSONDecodeError as e:
                LOGGER.error(f"Failed to parse JSON from Gemini API: {e}")
                await update.message.reply_text("Извините, произошла ошибка при обработке данных для презентации. Пожалуйста, попробуйте снова.")

        elif response_format == "text":
            LOGGER.info("Отправка запроса в Gemini API для генерации обычного текста.")
            payload = {
                "contents": contents,
                "generationConfig": {"temperature": 0.4}
            }
            gemini_response = await call_gemini_api(payload)
            clean_text = gemini_response.replace('<!DOCTYPE html>', '').replace('<html>', '').replace('<head>', '').replace('<body>', '').replace('</body>', '').replace('</html>', '').replace('<title>', '').replace('</title>', '').replace('</head>', '').replace('<div class="math-background">', '').replace('</div>', '').replace('<div class="default-background">', '').replace('<p>', '').replace('</p>', '').replace('<h1>', '').replace('</h1>', '')
            await update.message.reply_text(clean_text)

    except RetryAfter as e:
        LOGGER.warning(f"Flood control: Waiting for {e.retry_after} seconds.")
        await asyncio.sleep(e.retry_after)
        await update.message.reply_text("Извините, слишком много запросов. Пожалуйста, попробуйте снова через несколько секунд.")
    except NetworkError as e:
        LOGGER.error(f"Сетевая ошибка: {e}")
        await update.message.reply_text("Извините, произошла сетевая ошибка при обработке.")
    except Exception as e:
        LOGGER.error(f"Ошибка при обработке: {e}")
        await update.message.reply_text("Извините, произошла неизвестная ошибка.")

async def error_handler(update: object, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Обрабатывает ошибки в приложении и логирует их."""
    LOGGER.error("Произошла ошибка, но бот продолжит работу.")
    LOGGER.error(f"Update {update} caused error {context.error}")

def main() -> None:
    """Запускает бота."""
    if not BOT_TOKEN:
        LOGGER.error("Токен бота не найден в secrets.env. Пожалуйста, добавьте BOT_TOKEN.")
        exit(1)
        
    if not GEMINI_API_KEYS or GEMINI_API_KEYS == ['']:
        LOGGER.error("В файле secrets.env нет API-ключей. Пожалуйста, добавьте их.")
        exit(1)

    LOGGER.info(f"Найдено {len(GEMINI_API_KEYS)} API-ключей. Запуск бота...")
    
    application = Application.builder().token(BOT_TOKEN).build()
    
    # Команды и кнопки
    application.add_handler(CommandHandler("start", start_command_handler))
    application.add_handler(CommandHandler("reset", reset_command_handler))
    application.add_handler(CommandHandler("get_stars", get_stars_handler)) # НОВАЯ КОМАНДА
    application.add_handler(CommandHandler("donate", donate_command_handler)) # НОВАЯ КОМАНДА
    application.add_handler(CallbackQueryHandler(button_handler))
    
    # ОБРАБОТЧИКИ ДЛЯ ПЛАТЕЖЕЙ
    application.add_handler(PreCheckoutQueryHandler(pre_checkout_query_handler)) # НОВЫЙ ОБРАБОТЧИК
    application.add_handler(MessageHandler(filters.SUCCESSFUL_PAYMENT, successful_payment_callback)) # Уже был, но я его выделил
    
    # Универсальный обработчик для всех сообщений (текст, фото, файлы)
    application.add_handler(MessageHandler(filters.ALL & ~filters.COMMAND, handle_message))
    
    # Обработчик ошибок
    application.add_error_handler(error_handler)

    application.run_polling(allowed_updates=Update.ALL_TYPES)

if __name__ == "__main__":
    main()

